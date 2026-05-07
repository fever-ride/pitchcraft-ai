"""Generate starter .pptx templates for each project type.

Run once to populate backend/templates/pptx/ with minimal branded templates.
Each template defines slide layouts (title, content, section, blank) that ppt_builder uses.

NOTE: These are placeholder templates with colored header bars only.
Replace with designer-made .pptx files when available — ppt_builder picks them up by filename.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE_DIR = Path(__file__).parent.parent / "backend" / "templates" / "pptx"

THEMES = {
    "social": {
        "primary": RGBColor(0xFF, 0x6B, 0x6B),
        "secondary": RGBColor(0x4E, 0xCB, 0x71),
        "accent": RGBColor(0x33, 0x9A, 0xF0),
    },
    "pr": {
        "primary": RGBColor(0x2C, 0x3E, 0x50),
        "secondary": RGBColor(0x34, 0x98, 0xDB),
        "accent": RGBColor(0xE7, 0x4C, 0x3C),
    },
    "integrated": {
        "primary": RGBColor(0x1A, 0x1A, 0x2E),
        "secondary": RGBColor(0xE9, 0x4D, 0x6A),
        "accent": RGBColor(0x00, 0xB4, 0xD8),
    },
    "brand_refresh": {
        "primary": RGBColor(0x6C, 0x5C, 0xE7),
        "secondary": RGBColor(0xFD, 0xCB, 0x6E),
        "accent": RGBColor(0x00, 0xCE, 0xC9),
    },
}


def create_template(name: str, colors: dict) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add a single blank slide as placeholder so the template is valid
    blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Add a colored rectangle as header bar
    from pptx.shapes.autoshape import Shape
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colors["primary"]
    shape.line.fill.background()

    # Add template label
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(1.5))
    tf = txBox.text_frame
    tf.text = f"Pitchcraft — {name.replace('_', ' ').title()} Template"
    for para in tf.paragraphs:
        para.font.size = Pt(32)
        para.font.color.rgb = colors["primary"]
        para.alignment = PP_ALIGN.CENTER

    output = TEMPLATE_DIR / f"{name}.pptx"
    prs.save(str(output))
    print(f"  Created: {output}")


def main():
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    print("Generating PPTX templates...")
    for name, colors in THEMES.items():
        create_template(name, colors)
    # default is a copy of integrated
    create_template("default", THEMES["integrated"])
    print("Done.")


if __name__ == "__main__":
    main()
