import io

import pytest

pypdf = pytest.importorskip("pypdf")
docx = pytest.importorskip("docx")
pptx = pytest.importorskip("pptx")

from backend.core.rag.parser import parse_file, parse_file_structured


def test_parse_unsupported_extension():
    with pytest.raises(ValueError, match="Unsupported file type"):
        parse_file(b"content", "file.txt")


def test_parse_pdf_empty():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)

    result = parse_file(buf.read(), "test.pdf")
    assert result == ""


def test_parse_docx():
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello world")
    doc.add_paragraph("Second paragraph")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    result = parse_file(buf.read(), "test.docx")
    assert "Hello world" in result
    assert "Second paragraph" in result


def test_parse_pptx():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Title"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = parse_file(buf.read(), "deck.pptx")
    assert "Test Title" in result
    assert "[Slide 1]" in result


def test_parse_pptx_structured():
    from pptx import Presentation

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "First Slide"
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2.shapes.title.text = "Second Slide"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    doc = parse_file_structured(buf.read(), "deck.pptx")
    assert doc.total_pages == 2
    assert len(doc.segments) == 2
    assert doc.segments[0].slide_index == 1
    assert doc.segments[1].slide_index == 2
    assert "First Slide" in doc.segments[0].text


def test_parse_docx_structured():
    from docx import Document

    d = Document()
    d.add_paragraph("Para one")
    d.add_paragraph("Para two")
    buf = io.BytesIO()
    d.save(buf)
    buf.seek(0)

    doc = parse_file_structured(buf.read(), "test.docx")
    assert len(doc.segments) == 1
    assert "Para one" in doc.segments[0].text
    assert doc.segments[0].page_number is None
