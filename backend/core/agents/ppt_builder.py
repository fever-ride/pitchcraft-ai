"""PPT Builder: assembles slide content into a .pptx file using python-pptx."""
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE_DIR = Path(__file__).parent.parent.parent / "templates" / "pptx"
OUTPUT_DIR = Path(__file__).parent.parent.parent / "output"


def _get_template(project_type: str = "integrated") -> Path | None:
    template_path = TEMPLATE_DIR / f"{project_type}.pptx"
    if template_path.exists():
        return template_path
    default = TEMPLATE_DIR / "default.pptx"
    if default.exists():
        return default
    return None


def build_pptx(
    slides: list[dict],
    proposal_id: str,
    project_type: str = "integrated",
) -> str:
    """Build a .pptx file from slide content dicts. Returns the output file path."""
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    template = _get_template(project_type)
    if template:
        prs = Presentation(str(template))
    else:
        prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[-1]

    for slide_data in slides:
        content = slide_data.get("content", slide_data)
        slide = prs.slides.add_slide(blank_layout)

        title_text = content.get("title", "")
        body_text = content.get("body", "")
        bullets = content.get("bullets", [])

        left = Inches(0.8)
        top = Inches(0.6)
        width = Inches(11.5)
        height = Inches(1.2)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        for para in title_frame.paragraphs:
            para.font.size = Pt(28)
            para.font.bold = True

        if body_text:
            body_top = Inches(2.0)
            body_box = slide.shapes.add_textbox(left, body_top, width, Inches(1.0))
            body_frame = body_box.text_frame
            body_frame.text = body_text
            for para in body_frame.paragraphs:
                para.font.size = Pt(16)

        if bullets:
            bullet_top = Inches(3.2) if body_text else Inches(2.0)
            bullet_box = slide.shapes.add_textbox(left, bullet_top, width, Inches(4.0))
            bullet_frame = bullet_box.text_frame
            bullet_frame.text = ""
            for i, bullet in enumerate(bullets):
                if i == 0:
                    bullet_frame.paragraphs[0].text = f"• {bullet}"
                    bullet_frame.paragraphs[0].font.size = Pt(14)
                else:
                    p = bullet_frame.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(14)
                    p.space_before = Pt(6)

    output_path = OUTPUT_DIR / f"{proposal_id}.pptx"
    prs.save(str(output_path))
    return str(output_path)
