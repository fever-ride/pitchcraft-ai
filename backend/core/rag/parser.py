import io
import re
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path

# Lines appearing on more than this fraction of pages are treated as headers/footers.
_BOILERPLATE_THRESHOLD = 0.3

# Patterns that reliably indicate a page number line (no content value).
_PAGE_NUMBER_RE = re.compile(
    r"^[\-\s]*"              # optional leading dashes/spaces
    r"(第\s*\d+\s*[页P]"     # 第 3 页 / 第3P
    r"|[Pp]age\s*\d+"        # Page 3 / page 3
    r"|\d+\s*/\s*\d+"        # 3 / 20
    r"|\d+)"                 # bare page number
    r"[\-\s]*$"
)


@dataclass
class ParsedSegment:
    """A segment of text from a parsed document with source location metadata."""
    text: str
    page_number: int | None = None
    slide_index: int | None = None


@dataclass
class ParsedDocument:
    """Full parsed output with metadata for downstream contextual embedding."""
    segments: list[ParsedSegment] = field(default_factory=list)
    total_pages: int = 0

    @property
    def full_text(self) -> str:
        return "\n\n".join(seg.text for seg in self.segments if seg.text.strip())


def _detect_boilerplate(pages: list[str]) -> set[str]:
    """Return lines that appear on too many pages to carry real content.

    Counts normalised line occurrences across all pages.  Lines that appear on
    more than _BOILERPLATE_THRESHOLD of total pages, or match a page-number
    pattern, are considered boilerplate and should be stripped.
    """
    total = len(pages)
    if total < 2:
        return set()

    line_page_count: Counter = Counter()
    for page_text in pages:
        seen_on_page = {ln.strip() for ln in page_text.splitlines() if ln.strip()}
        line_page_count.update(seen_on_page)

    boilerplate: set[str] = set()
    for line, count in line_page_count.items():
        if count / total > _BOILERPLATE_THRESHOLD or _PAGE_NUMBER_RE.match(line):
            boilerplate.add(line)
    return boilerplate


def _strip_boilerplate(text: str, boilerplate: set[str]) -> str:
    """Remove boilerplate lines from a page's text."""
    if not boilerplate:
        return text
    lines = [ln for ln in text.splitlines() if ln.strip() not in boilerplate]
    return "\n".join(lines).strip()


def parse_pdf(file_bytes: bytes) -> ParsedDocument:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    raw_pages: list[tuple[int, str]] = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            raw_pages.append((i, text.strip()))

    boilerplate = _detect_boilerplate([t for _, t in raw_pages])

    segments = []
    for page_num, text in raw_pages:
        cleaned = _strip_boilerplate(text, boilerplate)
        if cleaned:
            segments.append(ParsedSegment(text=cleaned, page_number=page_num))
    return ParsedDocument(segments=segments, total_pages=len(reader.pages))


def parse_docx(file_bytes: bytes) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    full_text = "\n\n".join(paragraphs)
    segments = [ParsedSegment(text=full_text, page_number=None)]
    return ParsedDocument(segments=segments, total_pages=1)


def parse_pptx(file_bytes: bytes) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    segments = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            segments.append(ParsedSegment(text="\n".join(texts), slide_index=i))
    return ParsedDocument(segments=segments, total_pages=len(prs.slides))


def parse_file(file_bytes: bytes, filename: str) -> str:
    """Legacy interface: returns plain text with location markers.

    Use parse_file_structured() for the new pipeline with metadata tracking.
    """
    doc = parse_file_structured(file_bytes, filename)
    parts = []
    for seg in doc.segments:
        if not seg.text.strip():
            continue
        if seg.slide_index is not None:
            parts.append(f"[Slide {seg.slide_index}]\n{seg.text}")
        else:
            parts.append(seg.text)
    return "\n\n".join(parts)


def parse_file_structured(file_bytes: bytes, filename: str) -> ParsedDocument:
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(file_bytes)
    elif ext == ".docx":
        return parse_docx(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return parse_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
